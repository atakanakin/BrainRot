import os
import re
from PyPDF2 import PdfReader
import docx
from bs4 import BeautifulSoup
from pptx import Presentation

# WORDS_PER_MINUTE = 150  # avg words per minute
# MAX_SPEECH_DURATION_MINUTES = 2
# MAX_WORD_COUNT = WORDS_PER_MINUTE * MAX_SPEECH_DURATION_MINUTES


def clean_pdf_text(text: str) -> str:
    # """Clean extracted PDF text"""
    text = re.sub(r"/H\d+\s*", "", text)
    text = re.sub(r"\n(?=[a-z])", "", text)
    text = re.sub(r"[^\x20-\x7E\n\xA0-\xFF\u0100-\uFFFF]", "", text)
    text = re.sub(r"([a-z]{2,}[.!?])\s*(?=\d)", r"\1\n\n", text)
    text = re.sub(r"([a-z]{2,}[.!?])\s+(?=\d)", r"\1\n\n", text)
    return text.strip()


def pdf_to_text(input_path: str) -> str:
    reader = PdfReader(input_path, strict=False)
    extracted_text = []

    # Extract all pages first
    for page in reader.pages:
        # Use layout-preserving extraction
        page_text = page.extract_text()

        # Fix common PDF layout issues
        page_text = re.sub(r"(?<=\d)\s+(?=\d)", "", page_text)  # Fix split numbers
        page_text = re.sub(
            r"(?<=[a-z])-\s*\n\s*(?=[a-z])", "", page_text
        )  # Fix hyphenation

        # Remove headers/footers more accurately
        lines = page_text.split("\n")
        if len(lines) > 4:
            # Skip first 2 and last 2 lines if they look like headers/footers
            if re.match(r"^[\d\s©\-•]+$", lines[0]) or "Page" in lines[0]:
                lines = lines[2:-2]
            else:
                lines = lines[1:-1]

        # Fix bullet points and lists
        processed_lines = []
        for line in lines:
            # Fix bullet points
            line = re.sub(r"^[\s•⁃◦▪-]+\s*", "• ", line.strip())
            # Fix numbered lists
            line = re.sub(r"^\s*(\d+[\.)]) ", r"\1 ", line)
            processed_lines.append(line)

        # Join lines with proper spacing
        page_text = "\n".join(processed_lines)

        # Fix paragraph breaks
        page_text = re.sub(r"\n{3,}", "\n\n", page_text)  # Max 2 newlines
        page_text = re.sub(
            r"([.!?])\n(?=[A-Z])", r"\1\n\n", page_text
        )  # Add breaks after sentences

        # Fix math equations
        page_text = re.sub(r"(\d+)\s+(\d+)", r"\1\2", page_text)

        extracted_text.append(page_text)
    # Join pages with proper spacing
    text = "\n\n".join(extracted_text)

    # Final cleanup while preserving structure
    text = clean_pdf_text(text)

    return text


def extract_text(input_path, output_dir):
    """
    Extract text from a given file and save it as a .txt file.

    Args:
        input_path (str): Path to the input file.
        output_dir (str): Directory to save the output .txt file.

    Returns:
        str: Path to the generated .txt file.
    """

    file_name, file_extension = os.path.splitext(os.path.basename(input_path))
    file_extension = file_extension.lower()

    output_path = os.path.join(output_dir, f"{file_name}.txt")

    def save_text(content):
        with open(output_path, "w", encoding="utf-8") as file:
            file.write(content)
        # get the word count
        word_count = len(content.split(" "))
        return output_path, word_count

    if file_extension == ".txt":
        with open(input_path, "r", encoding="utf-8") as file:
            content = file.read()
        return save_text(content)

    elif file_extension == ".docx":
        doc = docx.Document(input_path)
        content = "\n".join([para.text for para in doc.paragraphs])
        return save_text(content)

    if file_extension.lower() == ".pdf":
        text = pdf_to_text(input_path)

        return save_text(text)

    elif file_extension == ".pptx":
        presentation = Presentation(input_path)
        content = []
        for slide in presentation.slides:
            slide_text = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    slide_text.append(shape.text)
            content.append("\n".join(slide_text))
        text = "\n".join(content)
        text = clean_pdf_text(text)
        text = re.sub(r"\n{3,}", "\n\n", text)
        return save_text(text)

    elif file_extension == ".html":
        with open(input_path, "r", encoding="utf-8") as file:
            soup = BeautifulSoup(file, "html.parser")
        text = soup.get_text(separator="\n")
        return save_text(text)

    else:
        raise ValueError(f"Unsupported file format: {file_extension}")
